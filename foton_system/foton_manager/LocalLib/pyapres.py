# Importações de bibliotecas necessárias
import json
import os
import time
import comtypes.client
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_ALIGN,MSO_VERTICAL_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement


# Classe para estilos de tabela
class TableStyle:
    # Constantes para diferentes estilos de tabela
    NoStyleNoGrid = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
    ThemedStyle1Accent1 = '{3C2FFA5D-87B4-456A-9821-1D50468CF0F}'
    ThemedStyle1Accent2 = '{284E427A-3D55-4303-BF80-6455036E1DE7}'
    ThemedStyle1Accent3 = '{69C7853C-536D-4A76-A0AE-DD22124D55A5}'
    ThemedStyle1Accent4 = '{775DCB02-9BB8-47FD-8907-85C794F793BA}'
    ThemedStyle1Accent5 = '{35758FB7-9AC5-4552-8A53-C91805E547FA}'
    ThemedStyle1Accent6 = '{08FB837D-C827-4EFA-A057-4D05807E0F7C}'
    NoStyleTableGrid = '{5940675A-B579-460E-94D1-54222C63F5DA}'
    ThemedStyle2Accent1 = '{D113A9D2-9D6B-4929-AA2D-F23B5EE8CBE7}'
    ThemedStyle2Accent2 = '{18603FDC-E32A-4AB5-989C-0864C3EAD2B8}'
    ThemedStyle2Accent3 = '{306799F8-075E-4A3A-A7F6-7FBC6576F1A4}'
    ThemedStyle2Accent4 = '{E269D01E-BC32-4049-B463-5C60D7B0CCD2}'
    ThemedStyle2Accent5 = '{327F97BB-C833-4FB7-BDE5-3F7075034690}'
    ThemedStyle2Accent6 = '{638B1855-1B75-4FBE-930C-398BA8C253C6}'
    LightStyle1 = '{9D7B26C5-4107-4FEC-AEDC-1716B250A1EF}'
    LightStyle1Accent1 = '{3B4B98B0-60AC-42C2-AFA5-B58CD77FA1E5}'
    LightStyle1Accent2 = '{0E3FDE45-AF77-4B5C-9715-49D594BDF05E}'
    LightStyle1Accent3 = '{C083E6E3-FA7D-4D7B-A595-EF9225AFEA82}'
    LightStyle1Accent4 = '{D27102A9-8310-4765-A935-A1911B00CA55}'
    LightStyle1Accent5 = '{5FD0F851-EC5A-4D38-B0AD-8093EC10F338}'
    LightStyle1Accent6 = '{68D230F3-CF80-4859-8CE7-A43EE81993B5}'
    LightStyle2 = '{7E9639D4-E3E2-4D34-9284-5A2195B3D0D7}'
    LightStyle2Accent1 = '{69012ECD-51FC-41F1-AA8D-1B2483CD663E}'
    LightStyle2Accent2 = '{72833802-FEF1-4C79-8D5D-14CF1EAF98D9}'
    LightStyle2Accent3 = '{F2DE63D5-997A-4646-A377-4702673A728D}'
    LightStyle2Accent4 = '{17292A2E-F333-43FB-9621-5CBBE7FDCDCB}'
    LightStyle2Accent5 = '{5A111915-BE36-4E01-A7E5-04B1672EAD32}'
    LightStyle2Accent6 = '{912C8C85-51F0-491E-9774-3900AFEF0FD7}'
    LightStyle3 = '{616DA210-FB5B-4158-B5E0-FEB733F419BA}'
    LightStyle3Accent1 = '{BC89EF96-8CEA-46FF-86C4-4CE0E7609802}'
    LightStyle3Accent2 = '{5DA37D80-6434-44D0-A028-1B22A696006F}'
    LightStyle3Accent3 = '{8799B23B-EC83-4686-B30A-512413B5E67A}'
    LightStyle3Accent4 = '{ED083AE6-46FA-4A59-8FB0-9F97EB10719F}'
    LightStyle3Accent5 = '{BDBED569-4797-4DF1-A0F4-6AAB3CD982D8}'
    LightStyle3Accent6 = '{E8B1032C-EA38-4F05-BA0D-38AFFFC7BED3}'
    MediumStyle1 = '{793D81CF-94F2-401A-BA57-92F5A7B2D0C5}'
    MediumStyle1Accent1 = '{B301B821-A1FF-4177-AEE7-76D212191A09}'
    MediumStyle1Accent2 = '{9DCAF9ED-07DC-4A11-8D7F-57B35C25682E}'
    MediumStyle1Accent3 = '{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}'
    MediumStyle1Accent4 = '{1E171933-4619-4E11-9A3F-F7608DF75F80}'
    MediumStyle1Accent5 = '{FABFCF23-3B69-468F-B69F-88F6DE6A72F2}'
    MediumStyle1Accent6 = '{10A1B5D5-9B99-4C35-A422-299274C87663}'
    MediumStyle2 = '{073A0DAA-6AF3-43AB-8588-CEC1D06C72B9}'
    MediumStyle2Accent1 = '{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}'
    MediumStyle2Accent2 = '{21E4AEA4-8DFA-4A89-87EB-49C32662AFE0}'
    MediumStyle2Accent3 = '{F5AB1C69-6EDB-4FF4-983F-18BD219EF322}'
    MediumStyle2Accent4 = '{00A15C55-8517-42AA-B614-E9B94910E393}'
    MediumStyle2Accent5 = '{7DF18680-E054-41AD-8BC1-D1AEF772440D}'
    MediumStyle2Accent6 = '{93296810-A885-4BE3-A3E7-6D5BEEA58F35}'
    MediumStyle3 = '{8EC20E35-A176-4012-BC5E-935CFFF8708E}'
    MediumStyle3Accent1 = '{6E25E649-3F16-4E02-A733-19D2CDBF48F0}'
    MediumStyle3Accent2 = '{85BE263C-DBD7-4A20-BB59-AAB30ACAA65A}'
    MediumStyle3Accent3 = '{EB344D84-9AFB-497E-A393-DC336BA19D2E}'
    MediumStyle3Accent4 = '{EB9631B5-78F2-41C9-869B-9F39066F8104}'
    MediumStyle3Accent5 = '{74C1A8A3-306A-4EB7-A6B1-4F7E0EB9C5D6}'
    MediumStyle3Accent6 = '{2A488322-F2BA-4B5B-9748-0D474271808F}'
    MediumStyle4 = '{D7AC3CCA-C797-4891-BE02-D94E43425B78}'
    MediumStyle4Accent1 = '{69CF1AB2-1976-4502-BF36-3FF5EA218861}'
    MediumStyle4Accent2 = '{8A107856-5554-42FB-B03E-39F5DBC370BA}'
    MediumStyle4Accent3 = '{0505E3EF-67EA-436B-97B2-0124C06EBD24}'
    MediumStyle4Accent4 = '{C4B1156A-380E-4F78-BDF5-A606A8083BF9}'
    MediumStyle4Accent5 = '{22838BEF-8BB2-4498-84A7-C5851F593DF1}'
    MediumStyle4Accent6 = '{16D9F66E-5EB9-4882-86FB-DCBF35E3C3E4}'
    DarkStyle1 = '{E8034E78-7F5D-4C2E-B375-FC64B27BC917}'
    DarkStyle1Accent1 = '{125E5076-3810-47DD-B79F-674D7AD40C01}'
    DarkStyle1Accent2 = '{37CE84F3-28C3-443E-9E96-99CF82512B78}'
    DarkStyle1Accent3 = '{D03447BB-5D67-496B-8E87-E561075AD55C}'
    DarkStyle1Accent4 = '{E929F9F4-4A8F-4326-A1B4-22849713DDAB}'
    DarkStyle1Accent5 = '{8FD4443E-F989-4FC4-A0C8-D5A2AF1F390B}'
    DarkStyle1Accent6 = '{AF606853-7671-496A-8E4F-DF71F8EC918B}'
    DarkStyle2 = '{5202B0CA-FC54-4496-8BCA-5EF66A818D29}'
    DarkStyle2Accent1Accent2 = '{0660B408-B3CF-4A94-85FC-2B1E0A45F4A2}'
    DarkStyle2Accent3Accent4 = '{91EBBBCC-DAD2-459C-BE2E-F6DE35CF9A28}'
    DarkStyle2Accent5Accent6 = '{46F890A9-2807-4EBB-B81D-B2AA78EC7F39}'


# Funções

def file_exists(file_path):
    """
    Verifica se um arquivo existe no caminho especificado.

    Args:
        file_path (str): O caminho do arquivo a ser verificado.

    Returns:
        bool: True se o arquivo existir, False caso contrário.
    """
    if not os.path.isfile(input_path):
        raise FileNotFoundError(f"Arquivo de entrada '{input_path}' não encontrado.")
    else:
        return os.path.isfile(file_path)

def open_presentation(input_file_path):
    """
    Abre uma apresentação PowerPoint a partir de um arquivo.

    Args:
        input_file_path (str): O caminho do arquivo da apresentação.

    Returns:
        Presentation: Um objeto Presentation representando a apresentação.
    
    Raises:
        FileNotFoundError: Se o arquivo especificado não for encontrado.
    """
    if file_exists(input_file_path):
        return Presentation(input_file_path)
    else:
        raise FileNotFoundError(f"Arquivo '{input_file_path}' não encontrado.")

def save_presentation(presentation, output_file_path):
    """
    Salva uma apresentação PowerPoint em um arquivo.

    Args:
        presentation (Presentation): O objeto Presentation da apresentação.
        output_file_path (str): O caminho do arquivo de saída.

    Returns:
        None
    
    Prints:
        str: Mensagem de erro se ocorrer algum problema ao salvar a apresentação.
    """
    try:
        presentation.save(output_file_path)
    except Exception as e:
        print(f"Erro ao salvar a apresentação: {str(e)}")

def ppt_to_pdf(input_file_path, output_file_path, format_type=32):
    """
    Converte uma apresentação PowerPoint em um arquivo PDF.

    Args:
        input_file_path (str): O caminho do arquivo de entrada (PPTX).
        output_file_path (str): O caminho do arquivo de saída (PDF).
        format_type (int): O formato de saída, padrão é 32 (PDF).

    Returns:
        None
    
    Prints:
        str: Mensagem de erro se ocorrer algum problema ao converter.
    """
    if not output_file_path.endswith('.pdf'):
        output_file_path += ".pdf"

    for i in range(5):
        if file_exists(input_file_path):
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1

            try:
                deck = powerpoint.Presentations.Open(input_file_path)
                deck.SaveAs(output_file_path, format_type)
                deck.Close()
                powerpoint.Quit()
                break
            except Exception as e:
                print(f"Erro ao processar o arquivo: {str(e)}")
        else:
            print(f"Arquivo '{input_file_path}' não encontrado. Tentando novamente em dois segundos.")
            time.sleep(2)
    else:
        print(f"Arquivo '{input_file_path}' não encontrado após 5 tentativas.")

def replace_text(presentation, replacements):
    """
    Substitui o texto em uma apresentação PowerPoint de acordo com um dicionário de substituições.

    Args:
        presentation (Presentation): O objeto Presentation da apresentação.
        replacements (dict): Um dicionário de substituições no formato {texto_antigo: texto_novo}.

    Returns:
        None
    """
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for old_text, new_text in replacements.items():
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)

def insert_styled_table(presentation, table_data, slide_index):
    """
    Insere uma tabela estilizada em uma apresentação PowerPoint.

    Args:
        presentation (Presentation): O objeto Presentation da apresentação.
        table_data (list): Uma lista de listas representando os dados da tabela.
        slide_index (int): O índice do slide onde a tabela será inserida.

    Returns:
        None
    
    Raises:
        ValueError: Se os dados da tabela não estiverem no formato correto.
    """
    # Verifica se os dados da tabela são válidos
    if not isinstance(table_data, list) or not all(isinstance(row, list) for row in table_data):
        raise ValueError("table_data must be a list of lists")
    if not all(len(row) == len(table_data[0]) for row in table_data):
        raise ValueError("All rows in table_data must have the same length")

    # Seleciona o slide
    slide = presentation.slides[slide_index]

    # Cria uma nova tabela
    x, y, cx, cy = Cm(1.45), Cm(4), Cm(32), Cm(12)
    tableFrame = slide.shapes.add_table(len(table_data), len(table_data[0]), x, y, cx, cy)
    
    table = tableFrame.table
  
    # Preenche a tabela com dados
    for row_num, row in enumerate(table_data):
        for col_num, cell_value in enumerate(row):
            cell = table.cell(row_num, col_num)
            cell.text = cell_value

    #insere o estilo da tablela
    set_table_style(tableFrame)

def set_table_style(graphicFrame):
    """
    Define o estilo de uma tabela em uma apresentação PowerPoint.

    Args:
        graphicFrame: O objeto que representa a tabela na apresentação.

    Returns:
        None
    """
    # Access the Table object
    tbl = graphicFrame._element.graphic.graphicData.tbl
    style_id = TableStyle.MediumStyle2Accent3
    tbl[0][-1].text = style_id

    # Set font size of first row to 14 points
    set_font_size(graphicFrame, 11)
    set_row_font_size(graphicFrame, 0, 14)
    set_bold_cell_rows(graphicFrame)

def set_font_size(graphicFrame, font_size):
    """
    Define o tamanho da fonte em células de uma tabela.

    Args:
        graphicFrame: O objeto que representa a tabela na apresentação.
        font_size (int): O tamanho da fonte a ser definido.

    Returns:
        None
    """
    tbl = graphicFrame.table

    for row in tbl.rows:
        for cell in row.cells:
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            if not cell.text_frame.paragraphs:
                p = cell.text_frame.add_paragraph()
                p.add_run()
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)

def set_row_font_size(graphicFrame, row_index, font_size):
    """
    Define o tamanho da fonte em uma linha específica de uma tabela.

    Args:
        graphicFrame: O objeto que representa a tabela na apresentação.
        row_index (int): O índice da linha a ser estilizada.
        font_size (int): O tamanho da fonte a ser definido.

    Returns:
        None
    """
    tbl = graphicFrame.table

    for cell in tbl.rows[row_index].cells:
        if not cell.text_frame.paragraphs:
            run.font.size = Pt(font_size)
            p = cell.text_frame.add_paragraph()
            p.add_run()
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)

def set_bold_cell_rows(graphicFrame):
    """
    Define células de uma linha específica como negrito em uma tabela.

    Args:
        graphicFrame: O objeto que representa a tabela na apresentação.

    Returns:
        None
    """
    # Access the Table object
    tbl = graphicFrame.table

    # Iterate over all rows in the table
    for row_index, row in enumerate(tbl.rows):
        # Count the number of non-empty cells in the row
        non_empty_cells = 0
        for cell in row.cells:
            if cell.text.strip() != "":
                non_empty_cells += 1

        # If the row contains only one non-empty cell or is the last row, set its text to bold
        if non_empty_cells == 1 or row_index == len(tbl.rows) - 1:
            for cell in row.cells:
                if not cell.text_frame.paragraphs:
                    p = cell.text_frame.add_paragraph()
                    p.add_run()
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True

def ppt_to_pdf(input_file_path, output_file_path, format_type=32):
    """
    Converte uma apresentação PowerPoint em um arquivo PDF.

    Args:
        input_file_path (str): O caminho do arquivo de entrada (PPTX).
        output_file_path (str): O caminho do arquivo de saída (PDF).
        format_type (int): O formato de saída, padrão é 32 (PDF).

    Returns:
        None
    
    Prints:
        str: Mensagem de erro se ocorrer algum problema ao converter.
    """
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    if output_file_path[-3:] != 'pdf':
        output_file_path = output_file_path + ".pdf"

    for i in range(5):
        if os.path.exists(input_file_path):
            deck = powerpoint.Presentations.Open(input_file_path)
            deck.SaveAs(output_file_path, format_type)
            deck.Close()
            powerpoint.Quit()
            break
        else:
            print(f"Arquivo '{input_file_path}' não encontrado. Tentando novamente em dois segundos.")
            time.sleep(2)
    else:
        print(f"Arquivo '{input_file_path}' não encontrado após 5 tentativas.")

def load_data_from_txt(txt_file_path):
    """
    Carrega dados de substituição de um arquivo de texto.

    Args:
        txt_file_path (str): O caminho do arquivo de texto.

    Returns:
        dict or None: Um dicionário de substituições ou None se o arquivo não for encontrado.
    """
    try:
        with open(txt_file_path, 'r') as txt_file:
            replacements = {}
            for line in txt_file:
                old_text, new_text = line.strip().split(';')
                replacements[old_text] = new_text
            return replacements
    except FileNotFoundError:
        return None

def load_data_from_json(json_file_path):
    """
    Carrega dados de uma tabela de um arquivo JSON.

    Args:
        json_file_path (str): O caminho do arquivo JSON.

    Returns:
        list or None: Uma lista de listas representando os dados da tabela ou None se o arquivo não for encontrado.
    """
    try:
        with open(json_file_path, 'r') as json_file:
            table_data = json.load(json_file)
            return table_data
    except FileNotFoundError:
        return None

def process_argument(arg):
    """
    Processa um argumento que pode ser um dicionário, uma lista ou um caminho para um arquivo.
    
    Args:
        arg: O argumento a ser processado.

    Returns:
        dict, list, or None: O argumento processado como um dicionário, lista ou None se não for um tipo suportado.
    """
    if isinstance(arg, dict) or isinstance(arg, list):
        return arg
    elif os.path.isfile(arg):
        # Se for um arquivo JSON, carrega os dados do arquivo
        if arg.lower().endswith('.json'):
            return load_data_from_json(arg)
        # Se for um arquivo de texto, carrega as substituições do arquivo
        elif arg.lower().endswith('.txt'):
            return load_data_from_txt(arg)
    return None

def edit_and_save_presentation(input_path, output_path, file_name, replacements=None, table_data=None, table_data_index=None):
    """
    Edita um arquivo de apresentação PowerPoint com opção de substituir texto e inserir uma tabela,
    e salva a edição em um novo arquivo PPTX. Também pode converter a apresentação para PDF.

    Args:
        input_path (str): O caminho do arquivo de entrada (PPTX).
        output_path (str): O caminho do diretório de saída.
        file_name (str): O nome do arquivo de saída (sem extensão).
        replacements (dict, optional): Um dicionário de substituições no formato {texto_antigo: texto_novo}.
        table_data (list, optional): Uma lista de listas representando os dados da tabela.

    Returns:
        None
    
    Raises:
        ValueError: Se o formato de saída não for suportado.
    """
    try:
        # Verifica se o arquivo de entrada existe
        if not os.path.isfile(input_path):
            raise FileNotFoundError(f"Arquivo de entrada '{input_path}' não encontrado.")

        # Abre a apresentação
        presentation = open_presentation(input_path)

        # Processa o argumento replacements, se fornecido
        if replacements is not None:
            replacements = process_argument(replacements)

        # Processa o argumento table_data, se fornecido
        if table_data is not None:
            table_data = process_argument(table_data)

        # Substitui o texto, se houver dados de substituição
        if replacements:
            replace_text(presentation, replacements)

        #Define o index de inserção da tabela, se não houver dados
        if not table_data_index:
            table_data_index = 7

        # Insere a tabela, se houver dados de tabela
        if table_data:
            insert_styled_table(presentation, table_data, table_data_index)

        # Salva a apresentação editada
        output_pptx = os.path.join(output_path, f"{file_name}.pptx")
        save_presentation(presentation, output_pptx)
        print(f"Apresentação editada salva em: '{output_pptx}'.")

        # Converte a apresentação para PDF, se necessário
        output_pdf = os.path.join(output_path, f"{file_name}.pdf")
        ppt_to_pdf(output_pptx, output_pdf, format_type=32)
        print(f"Apresentação PDF editada salva em: '{output_pdf}'.")
    except FileNotFoundError as e:
        print(f"Erro: {str(e)} Arquivo não encontrado.")
    except Exception as e:
        print(f"Ocorreu um erro: {str(e)}")

# Exemplo de uso:

if __name__ == "__main__":
    try:
        # Diretório e nomes de arquivos
        path  = os.path.dirname(os.path.abspath(__file__))
        input_pptx = "teste1.pptx"
        input_path = os.path.join(path, input_pptx)
        file_name = f"apresentacao_editada"

        # Dicionário de substituições de texto
        replacements = {
            "nomeProjeto": "Projeto Residencial",
            "localProjeto": "casa dos Sonhos",
            "nomeCliente": "Lucas Antonio",
            "empregoCliente": "Arquiteto",
            "estadoCivilCliente": "Casado",
            "estiloProjeto": "Projeto Moderno",
            "tamanhoProjeto":"250,00m²",
            "ambientesProjeto":"sala, quarto, banheiro",
            "anoProjeto":"2023",
            "levprojeto":"jan",
            "epprojeto": "fev",
            "pbProjeto":"mar",
            "alProjeto":"abr",
            "peProjeto":"mai",
            "inicioProjeto":"set",
            "valorProjeto":"10.000,00",
            "projArqEng": "R$1.000,00",
            "percArqEng": "10",
            "procLegais":"R$1.000,00",
            "percProcLegais": "10",
            "totalParcial":"R$1.000,00",
            "percTotalParcial": "10",
            "execInfra": "R$1.000,00",
            "percExecInfra": "10",
            "execCUB":"R$1.000,00",
            "percExecCUB":"10",
            "execPais":"R$1.000,00",
            "percExecPais":"10",
            "execMob":"R$1.000,00",
            "percExecMob":"10",
            "totalExec":"R$1.000,00",
            "totalInss":"R$1.000,00",
            "percINSS": "10",
            "totalGeral":"R$1.000,00",

        }

        # Dados da tabela
        table_data = [
            ["Serviço","Valor", "Percentual"],
            ["PLANEJAMENTO", " "," "],
            ["Projetos de Arquitetura e Engenharia", "R$2.000,00" , "10%"],
            ["Procedimentos Legais Municipais", "R$3.000,00", "10%"],
            ["Total parcial - Planejamento0","R$5.000,00", "20%"],
            ["EXECUÇÃO"," "," "],
            ["Infraestrutura","R$1.000,00","10%"],
            ["CUB - Custo Unitário Básico","R$1.000,00","10%"],
            ["Paisagismo","R$1.000,00","10%"],
            ["Mobiliário Fixo","R$1.000,00","10%"],
            ["Total Parcial - EXECUÇÃO","R$1.000,00","10%"],
            ["IMPOSTOS E CUSTOS INDIRETOS"," "," "],
            ["INSS","R$1.000,00","10%"],
            ["CUSTOS CARTORIAIS","R$1.000,00","10%"],
            ["EXPECTATIVA DE DESPESAS TOTAIS","R$1.000,00","100%"],
        ]

        edit_and_save_presentation(input_path, path, file_name, replacements, table_data)
            
    except Exception as e:
        print(f"Ocorreu um erro durante o processamento: {str(e)}")
    